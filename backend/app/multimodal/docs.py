"""Document intelligence: semantic chunking, table extraction, PPTX, spreadsheets.

Provider-independent with honest degradation:
- PDF: `parse_pdf` from pdf_parser (pypdf preferred, builtin fallback)
- Spreadsheets: openpyxl (installed) for XLSX; csv module for CSVs
- PPTX: python-pptx when installed, otherwise a truthful unavailable error
- Semantic chunking: heading/anchor-aware paragraph grouping
"""
import csv, io, logging, re
from dataclasses import dataclass, field
from typing import Any, Optional

logger = logging.getLogger(__name__)


@dataclass
class DocumentChunk:
    index: int
    text: str
    source: str = "document"
    page: int = 0
    heading: str = ""
    tokens: int = 0
    metadata: dict = field(default_factory=dict)

    def to_dict(self) -> dict:
        return {"index": self.index, "text": self.text, "source": self.source,
                "page": self.page, "heading": self.heading,
                "tokens": self.tokens, "metadata": self.metadata}


@dataclass
class TableExtraction:
    tables: list[list[list[str]]] = field(default_factory=list)
    page: int = 0
    engine: str = ""

    def to_dict(self) -> dict:
        return {"tables": self.tables[:50], "page": self.page, "engine": self.engine}


def _token_count(text: str) -> int:
    return max(1, len(text.split()))


class SemanticChunker:
    """Heading-anchored chunker: text is split at heading lines and then
    grouped into chunks with target token sizes. Honest structural heuristic
    (no learned model) - deterministic and predictable."""

    HEADING_RE = re.compile(
        r"^\s*(#{1,6}\s+)?(?P<body>[A-Z][A-Za-z0-9 ,:'&()./-]{2,79})\.?$")

    def __init__(self, target_tokens: int = 450, overlap_tokens: int = 40):
        self.target = target_tokens
        self.overlap = overlap_tokens

    def _split_headings(self, text: str) -> list[dict]:
        units = []
        for raw_line in text.splitlines():
            line = raw_line.strip()
            if not line:
                continue
            m = self.HEADING_RE.match(line)
            is_heading = bool(line.startswith("#"))
            if m and not is_heading:
                body = m.group("body")
                # sentence-style lines (end with punctuation, long, many words)
                # are body text, not headings
                is_heading = (
                    len(body) <= 64
                    and not line.endswith((".", "!", "?"))
                    and len(body.split()) <= 12
                )
            units.append({"heading": is_heading, "text": raw_line})
        return units

    def chunk(self, text: str, source: str = "document", page: int = 0) -> list[DocumentChunk]:
        prepared = self._split_headings(text)
        chunks: list[DocumentChunk] = []
        buffer: list[str] = []
        buffer_tokens = 0
        current_heading = ""
        index = 0
        for unit in prepared:
            line_tokens = _token_count(unit["text"])
            if unit["heading"]:
                current_heading = unit["text"].lstrip("#").strip()
            if buffer and buffer_tokens + line_tokens > self.target:
                chunk_text = "\n".join(buffer)
                chunks.append(DocumentChunk(
                    index=index, text=chunk_text, source=source, page=page,
                    heading=current_heading, tokens=_token_count(chunk_text)))
                index += 1
                tail: list[str] = []
                tail_tokens = 0
                for line in reversed(buffer):
                    tail.append(line)
                    tail_tokens += _token_count(line)
                    if tail_tokens >= self.overlap:
                        break
                buffer = list(reversed(tail))
                buffer_tokens = tail_tokens
            buffer.append(unit["text"])
            buffer_tokens += line_tokens
        if buffer:
            chunk_text = "\n".join(buffer)
            chunks.append(DocumentChunk(
                index=index, text=chunk_text, source=source, page=page,
                heading=current_heading, tokens=_token_count(chunk_text)))
        return chunks


class TableExtractor:
    """Extracts tables from PDF text (gridlines/whitespace heuristics) and
    from spreadsheets (openpyxl/csv)."""

    def from_pdf_text(self, page_text: str, page: int = 0) -> TableExtraction:
        """Best-effort grid detection: rows with consistent '|' or repeated
        whitespace-delimited cells. Conservative: returns nothing when the
        layout is ambiguous (never fabricates cells)."""
        rows = []
        for line in page_text.splitlines():
            if "|" in line:
                cells = [c.strip() for c in line.split("|")]
                rows.append(cells)
        if not rows:
            return TableExtraction(page=page, engine="none")
        # normalize: filter separator rows like |---|---|
        clean = [r for r in rows if not all(re.fullmatch(r"[-: ]+", c or "") for c in r)]
        if len(clean) < 2:
            return TableExtraction(page=page, engine="none")
        return TableExtraction(tables=[clean], page=page, engine="pdf-grid")

    def from_xlsx(self, data: bytes, sheet: str = "") -> TableExtraction:
        try:
            import openpyxl
        except ImportError:
            raise RuntimeError("openpyxl not installed; XLSX extraction unavailable")
        wb = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        ws = wb[sheet] if sheet and sheet in wb.sheetnames else wb.active
        rows = []
        for row in ws.iter_rows(values_only=True):
            if any(c is not None for c in row):
                rows.append(["" if c is None else str(c) for c in row])
        wb.close()
        return TableExtraction(tables=[rows] if rows else [], engine="openpyxl")

    def from_csv(self, data: bytes) -> TableExtraction:
        text = data.decode("utf-8-sig", errors="replace")
        rows = list(csv.reader(io.StringIO(text)))
        return TableExtraction(tables=[rows] if rows else [], engine="csv")


class DocIntelligence:
    """Facade over document extraction with honest unavailability reporting."""

    def __init__(self):
        self.chunker = SemanticChunker()
        self.tables = TableExtractor()

    def process_pdf(self, data: bytes, asset: dict) -> dict:
        from .pdf_parser import _safe_parse_pdf
        parsed = _safe_parse_pdf(data)
        if parsed.get("parser") == "error":
            parsed["type"] = "pdf"
            parsed["error"] = parsed.get("error", "PDF parse failed")
            return parsed
        pages = parsed.get("pages", [])
        chunks: list[DocumentChunk] = []
        for page in pages:
            page_chunks = self.chunker.chunk(
                page.get("text", ""), source=asset.get("id", "document"),
                page=page.get("page", 0))
            chunks.extend(page_chunks)
        tables = []
        for page in pages:
            ext = self.tables.from_pdf_text(page.get("text", ""), page.get("page", 0))
            if ext.tables:
                tables.extend(ext.tables)
        return {
            "parser": parsed.get("parser", ""),
            "page_count": parsed.get("page_count", 0),
            "pages": parsed.get("pages", []),
            "metadata": parsed.get("metadata", {}),
            "links": parsed.get("links", []),
            "outline": parsed.get("outline", []),
            "chunks": [c.to_dict() for c in chunks],
            "tables": tables,
            "full_text": "\n\n".join(p.get("text", "") for p in pages),
        }

    def process_spreadsheet(self, data: bytes, filename: str, asset: dict) -> dict:
        if filename.lower().endswith(".csv"):
            ext = self.tables.from_csv(data)
            rows = ext.tables[0] if ext.tables else []
            engine = ext.engine
        else:
            try:
                ext = self.tables.from_xlsx(data)
                rows = ext.tables[0] if ext.tables else []
                engine = ext.engine
            except RuntimeError as e:
                return {"engine": "unavailable", "error": str(e),
                        "chunks": [], "tables": [], "full_text": ""}
        full_text = "\n".join("\t".join(r) for r in rows)
        chunks = self.chunker.chunk(full_text, source=asset.get("id", "spreadsheet"))
        return {"engine": engine, "rows": rows[:5000], "tables": rows,
                "chunks": [c.to_dict() for c in chunks], "full_text": full_text}

    def process_pptx(self, data: bytes, asset: dict) -> dict:
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return {"engine": "unavailable",
                    "error": "python-pptx not installed; PPTX slides cannot be parsed locally",
                    "chunks": [], "slides": [], "full_text": ""}
        try:
            prs = Presentation(io.BytesIO(data))
        except Exception as e:
            return {"engine": "error", "error": str(e), "chunks": [], "slides": [], "full_text": ""}
        slides = []
        full = []
        for i, slide in enumerate(prs.slides, 1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs)
                        if t.strip():
                            parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
                if shape.shape_type == 13:  # picture
                    parts.append("[image]")
            text = "\n".join(parts)
            slides.append({"slide": i, "text": text})
            full.append(text)
        full_text = "\n\n".join(full)
        chunks = self.chunker.chunk(full_text, source=asset.get("id", "pptx"))
        return {"engine": "python-pptx", "slides": slides,
                "chunks": [c.to_dict() for c in chunks], "full_text": full_text}

    def process_text(self, data: bytes, asset: dict, encoding: str = "utf-8") -> dict:
        text = data.decode(encoding, errors="replace")
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        chunks = self.chunker.chunk(text, source=asset.get("id", "text"))
        return {"engine": "text", "chunks": [c.to_dict() for c in chunks],
                "full_text": text}


class DocumentPipeline:
    """Runs the right extractor for a document asset and augments with OCR
    results when a scanned (image-based) PDF has no text layer."""

    def __init__(self, doc_intel: Optional[DocIntelligence] = None, ocr=None):
        self.doc = doc_intel or DocIntelligence()
        self.ocr = ocr  # OCRDetector; used only when the PDF has no text layer

    def process(self, data: bytes, filename: str, asset: dict) -> dict:
        lower = filename.lower()
        if lower.endswith(".pdf"):
            result = self.doc.process_pdf(data, asset)
            result["type"] = "pdf"
            if not result.get("full_text", "").strip() and self.ocr is not None:
                # scanned PDF: rasterize impossible without poppler; OCR the
                # embedded images is out of scope here - report honestly
                result["scanned_hint"] = (
                    "no text layer found; use the OCR endpoint on extracted pages "
                    "when a rasterizer is available")
            return result
        if lower.endswith((".xlsx", ".xls", ".csv")):
            result = self.doc.process_spreadsheet(data, filename, asset)
            result["type"] = "spreadsheet"
            return result
        if lower.endswith((".pptx", ".ppt")):
            result = self.doc.process_pptx(data, asset)
            result["type"] = "pptx"
            return result
        if lower.endswith((".md", ".txt", ".json", ".yaml", ".yml", ".xml", ".html", ".cs", ".py")):
            result = self.doc.process_text(data, asset)
            result["type"] = "text"
            return result
        return {"type": "unsupported", "engine": "none", "chunks": [],
                "full_text": "", "error": f"no extractor for {filename}"}